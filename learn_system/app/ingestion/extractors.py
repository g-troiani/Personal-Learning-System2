"""Text extraction functions for various document formats."""

from pathlib import Path
from typing import Tuple


def extract_pdf(file_path: str) -> str:
    """
    Extracts text content from a PDF file.

    Args:
        file_path: Path to the PDF file

    Returns:
        Extracted text content as a single string
    """
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    text_parts = []

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text)

    return '\n\n'.join(text_parts)


def extract_docx(file_path: str) -> str:
    """
    Extracts text content from a DOCX file.

    Args:
        file_path: Path to the DOCX file

    Returns:
        Extracted text content as a single string
    """
    from docx import Document

    doc = Document(file_path)
    text_parts = []

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text_parts.append(paragraph.text)

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                if cell.text.strip():
                    row_text.append(cell.text.strip())
            if row_text:
                text_parts.append(' | '.join(row_text))

    return '\n\n'.join(text_parts)


def extract_markdown(file_path: str) -> str:
    """
    Reads and returns Markdown file content.

    Args:
        file_path: Path to the Markdown file

    Returns:
        File content as a string
    """
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()


def extract_text_file(file_path: str) -> str:
    """
    Reads and returns plain text file content.

    Args:
        file_path: Path to the text file

    Returns:
        File content as a string
    """
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()


def extract_pptx(file_path: str) -> str:
    """
    Extracts text content from a PowerPoint PPTX file.

    Extracts text from:
    - Slide titles
    - Text boxes and shapes
    - Tables
    - Speaker notes

    Args:
        file_path: Path to the PPTX file

    Returns:
        Extracted text content as a single string with slide separators
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    text_parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"=== Slide {slide_num} ==="]

        # Extract title if present
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()
            if title:
                slide_text.append(f"Title: {title}")

        # Extract body text from all shapes
        for shape in slide.shapes:
            # Skip the title shape (already extracted)
            if shape == slide.shapes.title:
                continue

            # Extract text from text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        slide_text.append(para_text)

            # Extract text from tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        slide_text.append(" | ".join(row_text))

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes_text = notes_frame.text.strip()
                if notes_text:
                    slide_text.append(f"[Speaker Notes]: {notes_text}")

        text_parts.append("\n".join(slide_text))

    return "\n\n".join(text_parts)


def extract_text(file_path: str) -> Tuple[str, str]:
    """
    Dispatches to appropriate extractor based on file extension.

    Args:
        file_path: Path to the document file

    Returns:
        Tuple of (extracted text content, content type)

    Raises:
        ValueError: If file extension is not supported
    """
    path = Path(file_path)
    extension = path.suffix.lower()

    extractors = {
        '.pdf': (extract_pdf, 'pdf'),
        '.docx': (extract_docx, 'docx'),
        '.doc': (extract_docx, 'docx'),  # Try docx extractor for .doc
        '.md': (extract_markdown, 'markdown'),
        '.markdown': (extract_markdown, 'markdown'),
        '.txt': (extract_text_file, 'text'),
        '.text': (extract_text_file, 'text'),
        '.pptx': (extract_pptx, 'pptx'),
        '.ppt': (extract_pptx, 'pptx'),  # Try pptx extractor for .ppt
    }

    if extension not in extractors:
        supported = ', '.join(extractors.keys())
        raise ValueError(f"Unsupported file type: {extension}. Supported types: {supported}")

    extractor_func, content_type = extractors[extension]
    content = extractor_func(file_path)

    return content, content_type


def get_file_metadata(file_path: str) -> dict:
    """
    Extracts metadata from a file.

    Args:
        file_path: Path to the file

    Returns:
        Dictionary containing file metadata
    """
    path = Path(file_path)
    stat = path.stat()

    return {
        'filename': path.name,
        'extension': path.suffix.lower(),
        'size_bytes': stat.st_size,
        'modified_timestamp': stat.st_mtime,
    }
